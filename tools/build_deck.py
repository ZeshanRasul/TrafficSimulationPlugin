"""Build the MegaGrant submission deck.

Six slides, figures taken from Docs/video-figures.md and the README so the deck
cannot drift from what the repo already claims.

Unlike plot_congestion.py this one is not stdlib-only:

    pip install python-pptx svglib pymupdf

    python tools/build_deck.py

svglib and pymupdf are only used to rasterise congestion_plot.svg for slide 5 —
PowerPoint will not take the SVG directly at a predictable size. Rendering it
here rather than committing a PNG keeps the plot single-sourced.
"""

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
PLOT = REPO / "Docs" / "congestion_plot.svg"
OUT = REPO / "Docs" / "TrafficSimulationPlugin-MegaGrant.pptx"

# Same reference palette the plot uses, so deck and figure agree.
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
PANEL = RGBColor(0xF4, 0xF3, 0xEF)

FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.85)             # left margin
CONTENT_W = W - L * 2


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return slide


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, colour=INK, bold=False, space_after=0,
         space_before=0, first=False, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)

    if line:
        p.line_spacing = line

    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = FONT
    return p


def rect(slide, x, y, w, h, colour):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def heading(slide, kicker, title, accent=BLUE):
    """Eyebrow, title, and a short accent rule — repeated on every slide."""
    tf = textbox(slide, L, Inches(0.62), CONTENT_W, Inches(0.3))
    para(tf, kicker.upper(), 12, MUTED, bold=True, first=True)

    tf = textbox(slide, L, Inches(0.95), CONTENT_W, Inches(0.6))
    para(tf, title, 32, INK, bold=True, first=True)

    rect(slide, L, Inches(1.66), Inches(1.0), Emu(27432), accent)


def table(slide, x, y, w, rows, col_w, header_colour=INK, row_h=Inches(0.42),
          aligns=None, emphasis=(), size=14):
    """A plain table: no banding, no rules, alignment set per column.

    `emphasis` names rows that read as summaries — filled and inked like the
    header rather than like a body row.
    """
    n_rows, n_cols = len(rows), len(rows[0])
    aligns = aligns or (["l"] + ["r"] * (n_cols - 1))

    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows)
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    for i, width in enumerate(col_w):
        tbl.columns[i].width = width

    for r, row in enumerate(rows):
        tbl.rows[r].height = row_h
        strong = (r == 0 or r in emphasis)

        for c, text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if strong else SURFACE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.margin_top = cell.margin_bottom = 0

            p = cell.text_frame.paragraphs[0]
            p.alignment = (PP_ALIGN.RIGHT if aligns[c] == "r"
                           else PP_ALIGN.LEFT)

            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.name = FONT
            run.font.bold = strong
            run.font.color.rgb = header_colour if strong else INK_2

    return tbl


def bullets(tf, items, size=16, gap=11, colour=INK_2, first=True):
    """Lead phrase in ink, rest in secondary — reads as a definition list."""
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.25

        run = p.add_run()
        run.text = lead
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = FONT

        if rest:
            run = p.add_run()
            run.text = rest
            run.font.size = Pt(size)
            run.font.color.rgb = colour
            run.font.name = FONT


# ---------------------------------------------------------------- slide 1

def slide_title(prs):
    slide = blank(prs)

    rect(slide, Emu(0), Emu(0), Inches(0.16), H, BLUE)

    tf = textbox(slide, L, Inches(2.05), CONTENT_W, Inches(0.4))
    para(tf, "EPIC MEGAGRANT SUBMISSION", 13, MUTED, bold=True, first=True)

    tf = textbox(slide, L, Inches(2.55), Inches(10.5), Inches(1.6))
    para(tf, "Traffic Simulation Plugin", 54, INK, bold=True, first=True,
         line=1.05)
    para(tf, "for Unreal Engine 5", 54, BLUE, bold=True, line=1.05)

    rect(slide, L, Inches(4.62), Inches(1.6), Emu(36576), ORANGE)

    tf = textbox(slide, L, Inches(4.95), Inches(9.6), Inches(1.2))
    para(tf, "Believable street traffic as ordinary Actors — spline road "
             "networks, signalised junctions with real right-of-way "
             "arbitration, and the tooling to see why the traffic is doing "
             "what it is doing.",
         18, INK_2, first=True, line=1.35)

    tf = textbox(slide, L, Inches(6.42), CONTENT_W, Inches(0.5))
    para(tf, "Zeshan Rasul", 15, INK, bold=True, first=True)
    para(tf, "Working prototype  ·  MIT licensed  ·  Unreal Engine 5.8",
         13, MUTED, space_before=3)


# ---------------------------------------------------------------- slide 2

def slide_why(prs):
    slide = blank(prs)
    heading(slide, "The gap", "Most projects aren't building CitySample")

    tf = textbox(slide, L, Inches(2.05), Inches(5.6), Inches(3.6))
    para(tf, "Unreal already ships a large-scale traffic system in "
             "CitySample, built on Mass Entity and ZoneGraph. It is "
             "impressive and it scales to a city.",
         16, INK_2, first=True, line=1.35)
    para(tf, "It is also sample content rather than a supported drop-in "
             "plugin, and using it means adopting an ECS architecture and "
             "learning its tooling.",
         16, INK_2, space_before=12, line=1.35)
    para(tf, "The far more common need is a few hundred believable vehicles "
             "on a handful of streets — configured by designers, inspectable "
             "when something looks wrong.",
         16, INK, space_before=12, line=1.35)

    rows = [
        ("", "This plugin", "MassTraffic"),
        ("Architecture", "Ordinary Actors", "Mass Entity (ECS)"),
        ("Configuration", "Details panel", "ECS traits"),
        ("Scale", "Hundreds", "Tens of thousands"),
        ("Learning curve", "Place an actor", "Learn Mass"),
        ("Form", "Plugin", "Sample content"),
    ]

    table(slide, Inches(7.0), Inches(2.05), Inches(5.5), rows,
          [Inches(1.55), Inches(1.85), Inches(2.1)], aligns=["l", "l", "l"])

    tf = textbox(slide, Inches(7.0), Inches(5.05), Inches(5.5), Inches(1.0))
    para(tf, "It is not a MassTraffic replacement.", 15, ORANGE, bold=True,
         first=True)
    para(tf, "If you need city-scale crowds, use Mass. This is the option "
             "for the case where traffic has to read convincingly and be "
             "reasoned about.",
         15, INK_2, space_before=5, line=1.3)


# ---------------------------------------------------------------- slide 3

def slide_how(prs):
    slide = blank(prs)
    heading(slide, "Architecture", "Three ideas do most of the work")

    col_w = Inches(3.75)
    gap = Inches(0.44)

    columns = [
        ("01", "One lane abstraction",
         "Roads and junctions both implement ITrafficLaneProvider, so a "
         "vehicle drives a junction connector with exactly the same code it "
         "uses on a road. Roads evaluate lanes from their spline; junctions "
         "interpolate baked samples. A vehicle never needs to know which it "
         "is on.",
         "Connectors are sampled by true arc length, so speed stays constant "
         "through a turn."),
        ("02", "Junction arbitration",
         "Conflicts are derived from actual path crossings, then resolved "
         "FIFO by arrival ticket. Entry is refused unless the exit lane has "
         "room — no blocking the box. Clearance is per-pair and "
         "positional: a turning vehicle blocks oncoming traffic only until "
         "it is past the meeting point.",
         "This cannot deadlock: the earliest movable ticket is always "
         "grantable."),
        ("03", "Car following",
         "Target speed is the minimum of the desired speed, a Gipps-style "
         "safe speed \u221a(v\u00b2 + 2\u00b7b\u00b7gap), and a time-headway "
         "comfort limit. The safe-speed term is the collision guarantee; the "
         "headway term is what makes queues look natural.",
         "Both account for the leader's speed, so a steady platoon is not "
         "slowed merely for being close."),
    ]

    for i, (num, title, body, note) in enumerate(columns):
        x = L + i * (col_w + gap)

        tf = textbox(slide, x, Inches(2.05), col_w, Inches(0.35))
        para(tf, num, 15, ORANGE, bold=True, first=True)

        tf = textbox(slide, x, Inches(2.45), col_w, Inches(0.5))
        para(tf, title, 19, INK, bold=True, first=True, line=1.15)

        tf = textbox(slide, x, Inches(3.15), col_w, Inches(2.2))
        para(tf, body, 14, INK_2, first=True, line=1.35)

        rect(slide, x, Inches(5.62), col_w, Inches(1.05), PANEL)

        tf = textbox(slide, x + Inches(0.2), Inches(5.78),
                     col_w - Inches(0.4), Inches(0.75))
        para(tf, note, 13, INK, first=True, line=1.3)

    tf = textbox(slide, L, Inches(6.92), CONTENT_W, Inches(0.35))
    para(tf, "11 automation tests cover the lane maths and junction "
             "arbitration, including a test that every queued vehicle "
             "eventually gets through.",
         13, MUTED, first=True)


# ---------------------------------------------------------------- slide 4

def slide_performance(prs):
    slide = blank(prs)
    heading(slide, "Performance", "500 vehicles, 16 signalised junctions, "
                                  "3.4 ms")

    tf = textbox(slide, L, Inches(2.05), Inches(6.1), Inches(0.9))
    para(tf, "Simulation time is measured directly, not inferred from frame "
             "time — a frame-rate cap makes frame time meaningless, because "
             "the engine simply idles to hit the deadline.",
         15, INK_2, first=True, line=1.35)

    rows = [
        ("Vehicles", "Simulation", "\u00b5s / vehicle", "Flow"),
        ("50", "0.18 ms", "3.6", "60%"),
        ("200", "0.94 ms", "4.7", "38%"),
        ("500", "3.36 ms", "6.7", "23%"),
    ]

    table(slide, L, Inches(3.2), Inches(6.1), rows,
          [Inches(1.5), Inches(1.6), Inches(1.6), Inches(1.4)],
          row_h=Inches(0.46))

    tf = textbox(slide, L, Inches(5.3), Inches(6.1), Inches(1.6))
    para(tf, "Flow is reported beside every timing because a saturated "
             "network is a different workload from a free-flowing one, and a "
             "timing figure without it is not interpretable. These runs "
             "deliberately loaded the network to saturation.",
         14, MUTED, first=True, line=1.35)

    # Scaling panel.
    rect(slide, Inches(7.35), Inches(2.05), Inches(5.15), Inches(4.6), PANEL)

    tf = textbox(slide, Inches(7.7), Inches(2.4), Inches(4.45), Inches(0.4))
    para(tf, "HOW IT SCALES", 12, MUTED, bold=True, first=True)

    tf = textbox(slide, Inches(7.7), Inches(2.9), Inches(4.45), Inches(0.6))
    para(tf, "T = 3.23 \u00b5s\u00b7N + 7.0\u00d710\u207b\u2076 ms\u00b7N\u00b2",
         21, BLUE, bold=True, first=True)

    tf = textbox(slide, Inches(7.7), Inches(3.6), Inches(4.45), Inches(2.7))
    para(tf, "Fits to within 1.5% across the measured range.",
         14, INK_2, first=True, line=1.3)

    bullets(textbox(slide, Inches(7.7), Inches(4.05), Inches(4.45),
                    Inches(2.3)), [
        ("Linear below ~460 vehicles", " — the per-vehicle update dominates."),
        ("Quadratic above it", " — the O(N\u00b2) neighbour search takes "
                               "over."),
        ("~10 ms at 1000 vehicles", " extrapolated."),
        ("A spatial grid is the known fix", " and is on the roadmap."),
    ], size=14, gap=9)


# ---------------------------------------------------------------- slide 5

def slide_congestion(prs, image):
    slide = blank(prs)
    heading(slide, "Emergent behaviour",
            "Congestion forms — and clears — on its own", ORANGE)

    tf = textbox(slide, L, Inches(2.0), Inches(11.6), Inches(0.4))
    para(tf, "One signal plan is swung to starve an axis across all 16 "
             "junctions, then biased back to work the backlog off.",
         15, INK_2, first=True, line=1.3)

    slide.shapes.add_picture(str(image), Inches(0.72), Inches(2.6),
                             width=Inches(7.6))

    x = Inches(8.55)
    rect(slide, x, Inches(2.6), Inches(3.95), Inches(4.28), PANEL)

    tf = textbox(slide, x + Inches(0.32), Inches(2.88), Inches(3.3),
                 Inches(0.3))
    para(tf, "WHAT THE RUN SHOWS", 12, MUTED, bold=True, first=True)

    stats = [
        ("68% \u2192 45%", "network flow, baseline to the deepest point of "
                           "the restriction"),
        ("63 of 100", "vehicles stopped at once at peak queue"),
        ("66%", "flow recovered, back to baseline"),
    ]

    y = Inches(3.3)

    for value, label in stats:
        tf = textbox(slide, x + Inches(0.32), y, Inches(3.3), Inches(0.36))
        para(tf, value, 23, ORANGE, bold=True, first=True)

        tf = textbox(slide, x + Inches(0.32), y + Inches(0.38), Inches(3.3),
                     Inches(0.5))
        para(tf, label, 13, INK_2, first=True, line=1.22)

        y += Inches(0.98)

    rect(slide, x + Inches(0.32), Inches(6.2), Inches(3.3), Emu(18288), RULE)

    tf = textbox(slide, x + Inches(0.32), Inches(6.38), Inches(3.3),
                 Inches(0.5))
    para(tf, "Flow keeps falling for 11 seconds after the restriction is "
             "lifted.",
         13, INK, bold=True, first=True, line=1.25)


# ---------------------------------------------------------------- slide 6

def slide_roadmap(prs):
    slide = blank(prs)
    heading(slide, "Where it goes next", "Known limits, and what closes them")

    tf = textbox(slide, L, Inches(2.05), Inches(5.6), Inches(0.35))
    para(tf, "KNOWN LIMITS", 12, MUTED, bold=True, first=True)

    bullets(textbox(slide, L, Inches(2.5), Inches(5.5), Inches(3.6)), [
        ("No rerouting.", " Vehicles pick a random successor at each "
                          "junction. A blocked route stays blocked."),
        ("No lane changing.", " One slow lorry holds up everything behind "
                              "it."),
        ("Hundreds, not thousands.", " Actor-based, so ~500 vehicles at "
                                     "3.4 ms."),
        ("O(N\u00b2) neighbour search.", " Deliberately unoptimised until it "
                                         "is the dominant cost."),
        ("No pedestrians or parking.", " Junction approaches are "
                                       "single-lane."),
    ], size=15, gap=10)

    tf = textbox(slide, Inches(7.0), Inches(2.05), Inches(5.5), Inches(0.35))
    para(tf, "ROADMAP", 12, MUTED, bold=True, first=True)

    items = [
        ("Congestion-aware routing",
         "Per-lane cost and successor choice that avoids blocked routes. "
         "The single largest improvement available."),
        ("Lane changing and overtaking",
         "Needed before routing is fully useful."),
        ("Spatial partitioning",
         "Lifts the practical ceiling past a thousand vehicles."),
        ("Editor tooling",
         "Junction placement and signal-phase authoring in the viewport."),
        ("Multi-lane approaches",
         "With turn-restricted lanes."),
    ]

    y = Inches(2.5)

    for i, (title, body) in enumerate(items):
        tf = textbox(slide, Inches(7.0), y, Inches(0.4), Inches(0.3))
        para(tf, str(i + 1), 15, ORANGE, bold=True, first=True)

        tf = textbox(slide, Inches(7.4), y, Inches(5.1), Inches(0.3))
        para(tf, title, 15, INK, bold=True, first=True)

        tf = textbox(slide, Inches(7.4), y + Inches(0.26), Inches(5.1),
                     Inches(0.4))
        para(tf, body, 13, INK_2, first=True, line=1.25)

        y += Inches(0.82)

    rect(slide, L, Inches(6.42), CONTENT_W, Emu(27432), RULE)

    tf = textbox(slide, L, Inches(6.65), Inches(11.6), Inches(0.7))
    para(tf, "Actor-based, configured in the details panel, and inspectable "
             "when something looks wrong \u2014 for the many projects that "
             "need believable street traffic rather than city-scale crowds.",
         15, INK, first=True, line=1.3)


def rasterise(svg_path, png_path, scale=2.5):
    """SVG -> PDF -> PNG. reportlab's direct PNG backend needs cairo, which is
    not reliably installable on Windows; its PDF backend is pure Python."""
    import pymupdf
    from reportlab.graphics import renderPDF
    from svglib.svglib import svg2rlg

    drawing = svg2rlg(str(svg_path))
    pdf = renderPDF.drawToString(drawing)

    with pymupdf.open(stream=pdf, filetype="pdf") as doc:
        doc[0].get_pixmap(matrix=pymupdf.Matrix(scale, scale)).save(
            str(png_path))


# ---------------------------------------------------------------- slide 7

# Every engineering line names the roadmap items it delivers, so a reviewer can
# check the budget against slide 6 rather than taking "development" on trust.
BUDGET = [
    ("Congestion-aware routing and lane changing", "Roadmap 1–2",
     17_000),
    ("Spatial partitioning, editor tooling, multi-lane", "Roadmap 3–5",
     11_000),
    ("Demo content and sample project", "Ships with the plugin", 4_000),
    ("Documentation and tutorials", "Ships with the plugin", 3_000),
    ("Fab submission and community launch", "Distribution", 2_000),
    ("Testing, infrastructure and software", "Ongoing", 2_000),
    ("Contingency", "", 1_000),
]


def slide_funding(prs):
    slide = blank(prs)

    total = sum(amount for _, _, amount in BUDGET)

    heading(slide, "The ask",
            f"${total:,}, mapped to the roadmap", ORANGE)

    rows = [("Allocation", "Delivers", "Amount", "Share")]

    for label, delivers, amount in BUDGET:
        share = amount / total * 100
        rows.append((
            label,
            delivers,
            f"${amount:,}",
            f"{share:.1f}%".replace(".0%", "%"),
        ))

    rows.append(("Total", "", f"${total:,}", "100%"))

    table(slide, L, Inches(2.05), CONTENT_W, rows,
          [Inches(5.4), Inches(2.6), Inches(1.9), Inches(1.7)],
          aligns=["l", "l", "r", "r"],
          row_h=Inches(0.4),
          emphasis={len(rows) - 1})

    band_y = Inches(6.05)
    rect(slide, L, band_y, CONTENT_W, Inches(0.95), PANEL)

    tf = textbox(slide, L + Inches(0.3), band_y + Inches(0.18),
                 CONTENT_W - Inches(0.6), Inches(0.6))
    para(tf, "Delivered over six months, in the roadmap order above. "
             "Everything the grant funds is released under the MIT licence "
             "— free to the Unreal community, with no revenue share and "
             "no exclusivity.",
         14, INK, first=True, line=1.3)


def main():
    if not PLOT.exists():
        raise SystemExit(
            f"missing {PLOT} — run tools/plot_congestion.py first")

    tmp = Path(tempfile.mkdtemp())
    image = tmp / "congestion.png"
    rasterise(PLOT, image)

    prs = new_deck()

    slide_title(prs)
    slide_why(prs)
    slide_how(prs)
    slide_performance(prs)
    slide_congestion(prs, image)
    slide_roadmap(prs)
    slide_funding(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
